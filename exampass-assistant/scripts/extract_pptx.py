"""Extract content from PowerPoint files."""

import os
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt


def extract_text_from_shape(shape) -> str:
    """Extract text from a shape, handling text frames."""
    if not shape.has_text_frame:
        return ""
    text_parts = []
    for para in shape.text_frame.paragraphs:
        para_text = ""
        for run in para.runs:
            para_text += run.text
        if para_text.strip():
            text_parts.append(para_text.strip())
    return "\n".join(text_parts)


def extract_table_content(shape) -> List[List[str]]:
    """Extract table as 2D list of strings."""
    table = shape.table
    rows = []
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


def extract_pptx(filepath: str) -> Dict[str, Any]:
    """
    Extract all content from a PPTX file.
    Returns: {
        'slides': [{'number': int, 'title': str, 'text': str, 'tables': [...], 'notes': str}],
        'total_slides': int
    }
    """
    prs = Presentation(filepath)
    slides = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_data = {
            'number': slide_idx,
            'title': '',
            'text': '',
            'tables': [],
            'image_count': 0,
            'notes': '',
        }

        texts = []
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
                slide_data['title'] = extract_text_from_shape(shape)
            elif shape.has_table:
                slide_data['tables'].append(extract_table_content(shape))
            elif shape.shape_type == 13:  # Picture
                slide_data['image_count'] += 1
            else:
                t = extract_text_from_shape(shape)
                if t:
                    texts.append(t)

        slide_data['text'] = '\n'.join(texts)

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_data['notes'] = slide.notes_slide.notes_text_frame.text.strip()

        slides.append(slide_data)

    return {
        'slides': slides,
        'total_slides': len(slides),
        'filename': os.path.basename(filepath),
    }


def build_text_summary(data: Dict[str, Any]) -> str:
    """Build a readable text summary from extracted PPTX data."""
    lines = [f"# {data['filename']}", f"共 {data['total_slides']} 张幻灯片\n"]
    for slide in data['slides']:
        lines.append(f"\n## 幻灯片 {slide['number']}")
        if slide['title']:
            lines.append(f"### {slide['title']}")
        if slide['text']:
            lines.append(slide['text'])
        for table in slide['tables']:
            lines.append("\n| " + " | ".join(table[0]) + " |")
            lines.append("|" + "|".join(["---"] * len(table[0])) + "|")
            for row in table[1:]:
                lines.append("| " + " | ".join(row) + " |")
        if slide['notes']:
            lines.append(f"\n> 备注: {slide['notes']}")
    return '\n'.join(lines)
