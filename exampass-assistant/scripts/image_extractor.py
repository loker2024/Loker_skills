"""Extract images from PPTX, DOCX, and PDF files."""

import os
import io
import zipfile
from pathlib import Path
from typing import List, Tuple
from PIL import Image

import pptx
import docx
import fitz  # pymupdf


def extract_from_pptx(pptx_path: str, output_dir: str) -> List[str]:
    """Extract all images from a PPTX file, save to output_dir, return paths."""
    paths = []
    prs = pptx.Presentation(pptx_path)
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    ext = image.content_type.split('/')[-1]
                    if ext == 'jpeg':
                        ext = 'jpg'
                    fname = f"slide{slide_idx+1}_img{shape_idx+1}.{ext}"
                    fpath = os.path.join(output_dir, fname)
                    with open(fpath, 'wb') as f:
                        f.write(image.blob)
                    paths.append(fpath)
                except Exception:
                    continue
    return paths


def extract_from_docx(docx_path: str, output_dir: str) -> List[str]:
    """Extract all images from a DOCX file via ZIP internals."""
    paths = []
    try:
        with zipfile.ZipFile(docx_path, 'r') as z:
            for name in z.namelist():
                if name.startswith('word/media/') and not name.endswith('/'):
                    img_data = z.read(name)
                    img_name = os.path.basename(name)
                    # Convert non-PNG images to PNG
                    try:
                        img = Image.open(io.BytesIO(img_data))
                        png_name = f"{os.path.splitext(img_name)[0]}.png"
                        fpath = os.path.join(output_dir, png_name)
                        img.save(fpath, 'PNG')
                        paths.append(fpath)
                    except Exception:
                        # fallback: save as-is
                        fpath = os.path.join(output_dir, img_name)
                        with open(fpath, 'wb') as f:
                            f.write(img_data)
                        paths.append(fpath)
    except Exception:
        pass
    return paths


def extract_from_pdf(pdf_path: str, output_dir: str) -> List[str]:
    """Extract all embedded images from a PDF file."""
    paths = []
    doc = fitz.open(pdf_path)
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        images = page.get_images(full=True)
        for img_idx, img_info in enumerate(images):
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                img_bytes = base_image["image"]
                ext = base_image["ext"]
                fname = f"page{page_idx+1}_img{img_idx+1}.{ext}"
                fpath = os.path.join(output_dir, fname)
                with open(fpath, 'wb') as f:
                    f.write(img_bytes)
                paths.append(fpath)
            except Exception:
                continue
    doc.close()
    return paths
