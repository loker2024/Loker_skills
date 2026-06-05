"""Shared fixtures for ExamPass tests."""

import os
import sys
import tempfile
import pytest

# Add scripts directory to path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'scripts'))

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DocxPt
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test outputs."""
    d = tempfile.mkdtemp(prefix="exampass_test_")
    yield d
    import shutil
    if os.path.exists(d):
        shutil.rmtree(d, ignore_errors=True)


@pytest.fixture
def sample_image_path(temp_dir):
    """Create a simple test image."""
    img = Image.new('RGB', (200, 100), color=(73, 109, 137))
    draw = ImageDraw.Draw(img)
    draw.rectangle((10, 10, 190, 90), outline=(255, 255, 255), width=2)
    draw.text((30, 40), "Test Image", fill=(255, 255, 255))
    path = os.path.join(temp_dir, "test_image.png")
    img.save(path)
    return path


def create_sample_pptx(path: str, include_image: bool = True, include_table: bool = True):
    """Create a sample PPTX with text, optional image and table."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "第一章 机器学习基础"
    if slide1.placeholders[1].has_text_frame:
        slide1.placeholders[1].text = "主讲：张教授"

    # Slide 2: Content with text
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "1.1 监督学习概述"
    body = slide2.placeholders[1]
    body.text = "监督学习是机器学习中最常见的范式。\n\n"
    body.text += "核心概念：\n"
    body.text += "- 训练集包含输入-输出对 (x, y)\n"
    body.text += "- 目标：学习映射函数 f: X → Y\n"
    body.text += "- 损失函数衡量预测误差：MSE = (1/n)Σ(yᵢ - ŷᵢ)²"

    # Slide 3: With table
    if include_table:
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        slide3.shapes.title.text = "常见算法对比"
        rows, cols = 4, 3
        table_shape = slide3.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(10), Inches(4))
        table = table_shape.table
        data = [
            ["算法", "类型", "适用场景"],
            ["线性回归", "回归", "连续值预测"],
            ["逻辑回归", "分类", "二分类问题"],
            ["SVM", "分类/回归", "高维数据"],
        ]
        for r in range(rows):
            for c in range(cols):
                table.cell(r, c).text = data[r][c]

    # Slide 4: With image
    if include_image:
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        slide4.shapes.title.text = "学习曲线示意图"
        # Create a simple test image
        img = Image.new('RGB', (400, 300), color=(200, 220, 240))
        draw = ImageDraw.Draw(img)
        draw.line([(50, 250), (150, 150), (250, 100), (350, 80)], fill=(0, 0, 200), width=3)
        draw.text((150, 50), "Training Curve", fill=(0, 0, 0))
        img_path = os.path.join(os.path.dirname(path), "_test_curve.png")
        img.save(img_path)
        slide4.shapes.add_picture(img_path, Inches(4), Inches(2.5), Inches(5), Inches(3.5))

    # Slide 5: Key points (empty slide content)
    if include_image is False and include_table is False:
        pass  # test with minimal content

    prs.save(path)


def create_sample_docx(path: str, include_image: bool = True, include_table: bool = True):
    """Create a sample DOCX with text, optional image and table."""
    doc = Document()

    doc.add_heading('第二章 深度学习基础', level=1)
    doc.add_paragraph('本章介绍深度学习的基本概念和核心算法。')

    doc.add_heading('2.1 神经网络结构', level=2)
    doc.add_paragraph(
        '人工神经网络（ANN）由多个层组成，每层包含若干神经元。'
        '输入层接收数据，隐藏层进行特征变换，输出层产生预测结果。'
        '激活函数引入非线性，使网络能够学习复杂模式。'
    )

    doc.add_heading('2.2 反向传播算法', level=2)
    doc.add_paragraph(
        '反向传播是训练神经网络的核心算法。它通过链式法则计算梯度，'
        '从输出层向输入层逐层传播误差。梯度下降更新权重以最小化损失函数。'
    )

    if include_table:
        doc.add_heading('激活函数对比', level=3)
        table = doc.add_table(rows=4, cols=3, style='Table Grid')
        data = [
            ["激活函数", "公式", "特点"],
            ["Sigmoid", "σ(x)=1/(1+e⁻ˣ)", "输出(0,1)，容易梯度消失"],
            ["ReLU", "f(x)=max(0,x)", "计算简单，缓解梯度消失"],
            ["Softmax", "S(xᵢ)=eˣⁱ/Σeˣʲ", "用于多分类输出层"],
        ]
        for r, row_data in enumerate(data):
            for c, cell_text in enumerate(row_data):
                table.cell(r, c).text = cell_text

    if include_image:
        img = Image.new('RGB', (400, 200), color=(240, 240, 240))
        draw = ImageDraw.Draw(img)
        for i in range(3):
            y = 50 + i * 60
            draw.ellipse([(80, y), (120, y+40)], fill=(100, 150, 200))
            draw.ellipse([(180, y), (220, y+40)], fill=(100, 150, 200))
            draw.ellipse([(280, y), (320, y+40)], fill=(100, 150, 200))
        img_path = os.path.join(os.path.dirname(path), "_test_nn.png")
        img.save(img_path)
        doc.add_picture(img_path, width=Inches(4))
        doc.add_paragraph('图：三层神经网络结构示意')

    doc.save(path)


def create_sample_pdf(path: str, include_table: bool = True):
    """Create a sample PDF with text and optional table."""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    c = canvas.Canvas(path, pagesize=A4)
    w, h = A4

    # Try to use a CJK font, fallback to Helvetica
    c.setFont("Helvetica", 16)
    c.drawString(50, h - 50, "Chapter 3: Convolutional Neural Networks")
    c.setFont("Helvetica", 12)
    c.drawString(50, h - 80, "3.1 CNN Architecture Overview")
    text_lines = [
        "Convolutional Neural Networks (CNNs) are specialized for processing grid-like data.",
        "",
        "Key Components:",
        "1. Convolutional Layer: Applies filters to extract features",
        "2. Pooling Layer: Reduces spatial dimensions (max pooling, average pooling)",
        "3. Fully Connected Layer: Final classification based on extracted features",
        "",
        "The convolution operation: output(i,j) = sum_m sum_n input(i+m, j+n) * kernel(m,n)",
        "",
        "Common CNN Architectures:",
        "- LeNet-5: Early digit recognition (1998)",
        "- AlexNet: ImageNet winner 2012, introduced ReLU and Dropout",
        "- VGG: Very deep with small 3x3 filters",
        "- ResNet: Introduced skip connections to train 152+ layers",
    ]
    y = h - 110
    for line in text_lines:
        c.drawString(50, y, line)
        y -= 18
        if y < 50:
            c.showPage()
            y = h - 50

    if include_table:
        c.showPage()
        c.setFont("Helvetica", 14)
        c.drawString(50, h - 50, "CNN Architecture Comparison")
        c.setFont("Helvetica", 10)
        table_data = [
            ["Architecture", "Year", "Layers", "Top-5 Error"],
            ["AlexNet", "2012", "8", "15.3%"],
            ["VGG-16", "2014", "16", "7.3%"],
            ["ResNet-152", "2015", "152", "4.5%"],
        ]
        y = h - 80
        for row in table_data:
            x = 50
            for cell in row:
                c.drawString(x, y, cell)
                x += 120
            y -= 20

    c.save()
